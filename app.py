import os
import streamlit as st
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.vectorstores import Chroma
from langchain.chains.question_answering import load_qa_chain
from langchain.llms import OpenAI
from langchain.callbacks import get_openai_callback
from PyPDF2 import PdfReader
from apikey import api_key
import docx2txt
import PyPDF2
import textract
from pptx import Presentation
from io import BytesIO
from dotenv import load_dotenv
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.embeddings import SentenceTransformerEmbeddings
from transformers import AutoModel, AutoTokenizer
import numpy as np

def extract_text_from_docx(docx_bytes):
    return docx2txt.process(BytesIO(docx_bytes))

def extract_text_from_pdf(pdf_bytes):
    pdf_text = ""
    pdf_reader = PyPDF2.PdfReader(BytesIO(pdf_bytes))
    for page in pdf_reader.pages:
        pdf_text += page.extract_text()
    return pdf_text

def extract_text_from_ppt(ppt_bytes):
    ppt_text = ""
    presentation = Presentation(BytesIO(ppt_bytes))
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                ppt_text += shape.text + "\n"
    return ppt_text

def process_uploaded_files(docs):
    all_text = ""
    for doc in docs:
        if doc.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc_text = extract_text_from_docx(doc.read())
        elif doc.type == "application/pdf":
            doc_text = extract_text_from_pdf(doc.read())
        elif doc.type == "application/vnd.ms-powerpoint":
            doc_text = extract_text_from_ppt(doc.read())
        else:
            doc_text = textract.process(doc.read()).decode("utf-8", errors="ignore")
        all_text += doc_text + "\n"  # Concatenate text
        #all_text.append({"page_content": doc_text})  # Create objects with page_content attribute
    return all_text


def get_vectorstore(chunks):
    embeddings = SentenceTransformerEmbeddings(model_name="all-MiniLM-L6-v2")

    if not chunks:
        return None
     
    try:
        vectorstore = Chroma.from_documents(chunks, embeddings)
        return vectorstore
    
    except Exception as e:
        return None  


def main():
    load_dotenv()
    os.environ['OPENAI_API_KEY'] = 'Your API key here'
    st.set_page_config(page_title="Query your PDFs", page_icon=":scroll:")
    st.header("The ultimate PDF whisperer 💬")
    
    # upload files
    pdfs = st.file_uploader("Upload your PDFs", type=["docx", "pdf", "ppt", "txt"], accept_multiple_files=True)
    
    # process each uploaded PDF
    if pdfs is not None:
        # extract the text
        text = process_uploaded_files(pdfs)

        # split into chunks
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )
        chunks = text_splitter.create_documents(text)

        if chunks is not None:        
            knowledge_base = get_vectorstore(chunks)

        user_question = st.text_input(f"Ask a question about PDF:")

        
        
        #convert the user question to an embed 
        docs = knowledge_base.similarity_search(user_question)

        

        llm = OpenAI()
        chain = load_qa_chain(llm, chain_type="stuff")
        with get_openai_callback() as cb:
            response = chain.run(input_documents=docs, question=user_question)
            print(cb)
        
        st.write(response)

if __name__ == '__main__':
    main()
